import numpy as np
import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from datetime import datetime, timedelta


def defect_pareto(engine, top_n=15):
    """Main pareto analysis function - updated for database integration"""
    st.markdown("## 📈 Quality Intelligence Dashboard")
    
    # Load data from database
    df = load_data_from_db(engine)
    
    if df.empty:
        st.info("No data available for analysis")
        return
    
    # Four tabs for different analyses
    tab1, tab2, tab3, tab4 = st.tabs([
        "🔧 Chronic Issues", 
        "👥 Operator Trends", 
        "📊 Daily Performance",
        "🔍 Advanced Analysis"
    ])
    
    with tab1:
        render_chronic_issues(engine, top_n)
    
    with tab2:
        render_operator_trends(engine)
    
    with tab3:
        render_performance_trends(engine)
    
   # with tab4:
    #    render_advanced_analysis(engine)

def run_query(query, engine):
    """Run SQL query and return DataFrame"""
    if not hasattr(engine, "connect"):
        st.error(f"❌ Invalid engine passed to run_query (got {type(engine)})")
        return pd.DataFrame()
    try:
        return pd.read_sql(query, engine)
    except Exception as e:
        st.error(f"Query failed: {e}")
        return pd.DataFrame()

def load_data_from_db(engine):
    """Load data from database"""
    try:
        # Quick test to confirm engine works
        with engine.connect() as conn:
            pass
    except Exception as e:
        st.error(f"❌ Database engine test failed: {e}")
        return pd.DataFrame()

    query = """
        SELECT *
        FROM quality.clean_quality_data
        WHERE date >= CURRENT_DATE - INTERVAL '90 days'
    """
    df = run_query(query, engine)

    if df.empty:
        return df

    df['date'] = pd.to_datetime(df['date'], errors='coerce')
    df['date_only'] = df['date'].dt.date

    if 'id' not in df.columns:
        df['id'] = range(1, len(df) + 1)

    return df

import io
import zipfile
import pandas as pd
import numpy as np
import streamlit as st
import plotly.graph_objects as go

# Optional libs (use fallbacks if not installed)
try:
    import xlsxwriter  # for embedding images into Excel (preferred)
    HAS_XLSXWRITER = True
except Exception:
    HAS_XLSXWRITER = False

try:
    import openpyxl  # alternative to xlsxwriter for Excel; used for image insertion if available
    from openpyxl.drawing.image import Image as OpenpyxlImage
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except Exception:
    Presentation = None
    HAS_PPTX = False

# Pillow is required by openpyxl Image insertion from BytesIO in some environments
try:
    from PIL import Image as PILImage
    HAS_PIL = True
except Exception:
    HAS_PIL = False


def _make_png_bytes(fig, width=1200, height=600):
    """Return PNG bytes of a plotly figure where possible (kaleido or plotly write_image)."""
    try:
        # prefers kaleido-backed to_image
        png = fig.to_image(format="png", width=width, height=height)
        return png
    except Exception:
        try:
            buf = io.BytesIO()
            fig.write_image(buf, format="png", width=width, height=height)
            return buf.getvalue()
        except Exception:
            return None


def _create_excel_with_image(df, png_bytes):
    """Return bytes of an Excel workbook. Try xlsxwriter first, then openpyxl.
       If neither image-capable route is available, return a simple Excel bytes (data only) if possible.
    """
    out = io.BytesIO()

    # Try xlsxwriter (easiest for embedding image)
    if HAS_XLSXWRITER and png_bytes:
        try:
            with pd.ExcelWriter(out, engine="xlsxwriter") as writer:
                # summary
                summary_df = pd.DataFrame({
                    "metric": ["Total_Defects_in_selection", "Total_Scrap_in_selection", "Overall_Scrap_Rate_%", "Rows_returned"],
                    "value": [
                        int(df['defect_count'].sum()),
                        int(df['scrap_count'].sum()),
                        float((df['scrap_count'].sum() / df['defect_count'].sum() * 100) if df['defect_count'].sum() > 0 else 0.0),
                        len(df)
                    ]
                })
                summary_df.to_excel(writer, sheet_name="Summary", index=False)
                df.to_excel(writer, sheet_name="Top_Defects", index=False)
                # insert image
                workbook = writer.book
                worksheet = writer.sheets["Top_Defects"]
                img_buf = io.BytesIO(png_bytes)
                # place image to the right of the table
                worksheet.insert_image("H2", "chart.png", {"image_data": img_buf, "x_scale": 0.8, "y_scale": 0.8})
            return out.getvalue()
        except Exception:
            # fall through to other methods
            pass

    # Try openpyxl path with image insertion (requires openpyxl and pillow in many setups)
    if HAS_OPENPYXL and png_bytes and HAS_PIL:
        try:
            # First write sheets via pandas (openpyxl engine)
            with pd.ExcelWriter(out, engine="openpyxl") as writer:
                summary_df = pd.DataFrame({
                    "metric": ["Total_Defects_in_selection", "Total_Scrap_in_selection", "Overall_Scrap_Rate_%", "Rows_returned"],
                    "value": [
                        int(df['defect_count'].sum()),
                        int(df['scrap_count'].sum()),
                        float((df['scrap_count'].sum() / df['defect_count'].sum() * 100) if df['defect_count'].sum() > 0 else 0.0),
                        len(df)
                    ]
                })
                summary_df.to_excel(writer, sheet_name="Summary", index=False)
                df.to_excel(writer, sheet_name="Top_Defects", index=False)

            # Insert image using openpyxl directly
            out.seek(0)
            wb = openpyxl.load_workbook(out)
            ws = wb["Top_Defects"]
            img = OpenpyxlImage(io.BytesIO(png_bytes))
            # scale image so it fits reasonably
            img.width = int(img.width * 0.75)
            img.height = int(img.height * 0.75)
            ws.add_image(img, "H2")
            final = io.BytesIO()
            wb.save(final)
            return final.getvalue()
        except Exception:
            pass

    # Fallback: create a data-only Excel (no image) using available engines
    try:
        out2 = io.BytesIO()
        # prefer openpyxl if installed, otherwise let pandas choose default
        engine = "openpyxl" if HAS_OPENPYXL else None
        with pd.ExcelWriter(out2, engine=engine) as writer:
            summary_df = pd.DataFrame({
                "metric": ["Total_Defects_in_selection", "Total_Scrap_in_selection", "Overall_Scrap_Rate_%", "Rows_returned"],
                "value": [
                    int(df['defect_count'].sum()),
                    int(df['scrap_count'].sum()),
                    float((df['scrap_count'].sum() / df['defect_count'].sum() * 100) if df['defect_count'].sum() > 0 else 0.0),
                    len(df)
                ]
            })
            summary_df.to_excel(writer, sheet_name="Summary", index=False)
            df.to_excel(writer, sheet_name="Top_Defects", index=False)
        return out2.getvalue()
    except Exception:
        return None


def _create_pptx_with_image(df, png_bytes):
    """Return PPTX bytes if python-pptx is installed and png_bytes available, else None."""
    if not HAS_PPTX or not png_bytes:
        return None
    try:
        prs = Presentation()
        # Use a simple title-and-content layout if available
        layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        try:
            title = slide.shapes.title
            if title:
                title.text = "Chronic Quality Issues - Summary"
        except Exception:
            pass
        img_buf = io.BytesIO(png_bytes)
        # Add image (scale width to 9 inches)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(1.2), width=Inches(9.0))

        # Add a slide with top defects text summary
        slide2 = prs.slides.add_slide(layout)
        try:
            title2 = slide2.shapes.title
            if title2:
                title2.text = "Top Defects (selection)"
        except Exception:
            pass

        top5 = df.head(10)[["defect", "defect_count", "scrap_count", "scrap_rate"]]
        tx_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.0))
        tf = tx_box.text_frame
        for _, row in top5.iterrows():
            p = tf.add_paragraph()
            p.text = f"{row['defect']} — count: {int(row['defect_count']):,}, scrap: {int(row['scrap_count']):,}, rate: {row['scrap_rate']:.1f}%"

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception:
        return None


def _make_zip_bundle(files_dict):
    """Create an in-memory zip with provided filename->bytes mapping and return bytes."""
    out = io.BytesIO()
    with zipfile.ZipFile(out, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name, content in files_dict.items():
            if content is None:
                continue
            zf.writestr(name, content)
    return out.getvalue()


def render_chronic_issues(engine, top_n=15, debug=False, sort_by='scrap_rate'):
    """Main UI function: renders chart and produces download artifacts with robust fallbacks."""
    st.markdown("## 🔧 Chronic Quality Issues Analysis")
    st.info("Identifies defects with highest scrap rate impact for priority attention.")

    try:
        query = f"""
        SELECT 
            code_description AS defect,
            COUNT(*) AS defect_count,
            COUNT(CASE WHEN disposition = 'SCRAP' THEN 1 END) AS scrap_count
        FROM quality.clean_quality_data
        WHERE code_description IS NOT NULL AND code_description != ''
        GROUP BY code_description
        HAVING COUNT(*) > 0
        ORDER BY defect_count DESC
        LIMIT {top_n}
        """
        df = pd.read_sql(query, engine)

        if df.empty:
            st.warning("⚠️ No defect data available in database.")
            return

        df["defect_count"] = pd.to_numeric(df["defect_count"], errors="coerce").fillna(0).astype(int)
        df["scrap_count"] = pd.to_numeric(df["scrap_count"], errors="coerce").fillna(0).astype(int)

        df["scrap_rate"] = np.where(
            df["defect_count"] > 0,
            (df["scrap_count"] / df["defect_count"] * 100),
            0.0
        )
        df["scrap_rate"] = df["scrap_rate"].round(1)
        df["defect_percentage"] = (df["defect_count"] / df["defect_count"].sum() * 100).round(1)

        if sort_by == 'scrap_rate':
            df = df.sort_values("scrap_rate", ascending=False)
        else:
            df = df.sort_values("defect_count", ascending=False)

        if debug:
            st.write("DEBUG - df (used for chart)", df[['defect','defect_count','scrap_count','scrap_rate','defect_percentage']])
            st.write("DEBUG - sums", {
                "scrap_sum": int(df['scrap_count'].sum()),
                "defect_sum": int(df['defect_count'].sum()),
                "overall_rate": float((df['scrap_count'].sum() / df['defect_count'].sum() * 100) if df['defect_count'].sum() > 0 else 0)
            })

        # Build figure (same as before)
        fig = go.Figure()
        def get_scrap_rate_color(scrap_rate):
            if scrap_rate >= 70:
                return '#FF4444'
            elif scrap_rate >= 40:
                return '#FFAA44'
            elif scrap_rate >= 20:
                return '#44AAFF'
            else:
                return '#44FF88'

        colors = [get_scrap_rate_color(rate) for rate in df["scrap_rate"]]
        x_vals = list(df['scrap_rate'])
        y_vals = list(df['defect'])
        customdata = np.stack([df['defect_count'].astype(int),
                               df['scrap_count'].astype(int),
                               df['defect_percentage'].astype(float)], axis=1)

        fig.add_trace(go.Bar(
            x=x_vals,
            y=y_vals,
            orientation='h',
            marker=dict(color=colors, line=dict(color='darkgray', width=1)),
            hovertemplate=(
                "<b>%{y}</b><br>"
                "🔄 Scrap Rate: <b>%{x:.1f}%</b><br>"
                "📊 Total Defects: %{customdata[0]:,}<br>"
                "🗑️ Scrap Count: %{customdata[1]:,}<br>"
                "📈 Frequency: %{customdata[2]:.1f}%<br>"
                "<extra></extra>"
            ),
            customdata=customdata,
            name="Scrap Rate %"
        ))

        fig.update_layout(
            title=dict(text="🔥 Defects by Scrap Rate Impact (Highest to Lowest)", x=0.5, font=dict(size=18)),
            xaxis=dict(title="Scrap Rate (%)", range=[0, 100], showgrid=True, gridcolor='lightgray'),
            yaxis=dict(title="Defect Types", autorange="reversed"),
            height=max(500, len(df) * 35),
            template="plotly_white",
            margin=dict(l=180, r=20, t=80, b=50),
            hovermode="y unified",
            showlegend=False
        )

        fig.add_vline(x=70, line_dash="dash", line_color="red", annotation_text="Critical >70%")
        fig.add_vline(x=40, line_dash="dash", line_color="orange", annotation_text="High >40%")
        fig.add_vline(x=20, line_dash="dash", line_color="blue", annotation_text="Medium >20%")

        st.plotly_chart(fig, use_container_width=True)

        # -----------------------
        # Robust report creation & downloads
        # -----------------------
        st.markdown("### 📥 Downloadable Reports")

        # Build simple summary and CSV
        summary = {
            "Total_Defects_in_selection": int(df['defect_count'].sum()),
            "Total_Scrap_in_selection": int(df['scrap_count'].sum()),
            "Overall_Scrap_Rate_%": float((df['scrap_count'].sum() / df['defect_count'].sum() * 100) if df['defect_count'].sum() > 0 else 0.0),
            "Rows_returned": int(len(df))
        }
        csv_bytes = df.to_csv(index=False).encode('utf-8')

        # PNG of the chart (try to create; may require kaleido)
        png_bytes = _make_png_bytes(fig, width=1200, height=max(400, len(df) * 35))

        # Excel bytes (try xlsxwriter -> openpyxl -> data-only)
        excel_bytes = _create_excel_with_image(df, png_bytes)

        # PPTX bytes (if python-pptx + png available)
        pptx_bytes = _create_pptx_with_image(df, png_bytes)

        # If both excel and pptx are missing, create a ZIP bundle with CSV + PNG (if any)
        files = {
            "top_defects.csv": csv_bytes,
            "top_defects_chart.png": png_bytes,
            "top_defects_report.xlsx": excel_bytes,
            "top_defects_summary.pptx": pptx_bytes
        }
        zip_bytes = None
        # If at least one of excel/pptx/png exists, also produce zip for convenience
        if any(v is not None for v in [excel_bytes, pptx_bytes, png_bytes, csv_bytes]):
            zip_bytes = _make_zip_bundle(files)

        # Display helpful notes about available exports and missing deps
        dep_msgs = []
        if not HAS_XLSXWRITER and not HAS_OPENPYXL:
            dep_msgs.append("Excel image embedding unavailable (install XlsxWriter or openpyxl + pillow to embed chart).")
        if not png_bytes:
            dep_msgs.append("PNG export unavailable (install 'kaleido' or enable plotly image export).")
        if not HAS_PPTX:
            dep_msgs.append("PPTX creation unavailable (install 'python-pptx').")
        if dep_msgs:
            for m in dep_msgs:
                st.info(m)

        # Render download buttons in columns
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            st.download_button("CSV - Top Defects", data=csv_bytes, file_name="top_defects.csv", mime="text/csv")
        with c2:
            if excel_bytes:
                st.download_button("Excel (report)", data=excel_bytes, file_name="top_defects_report.xlsx",
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            else:
                st.button("Excel (report) - unavailable")
        with c3:
            if png_bytes:
                st.download_button("PNG Chart", data=png_bytes, file_name="top_defects_chart.png", mime="image/png")
            else:
                st.button("PNG Chart - unavailable")
        with c4:
            if pptx_bytes:
                st.download_button("PPTX (presentation)", data=pptx_bytes, file_name="top_defects_summary.pptx",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            else:
                st.button("PPTX - unavailable")

        # Provide a ZIP with everything that was produced for easy sharing
        if zip_bytes:
            st.download_button("ZIP - All artifacts", data=zip_bytes, file_name="top_defects_bundle.zip", mime="application/zip")

    except Exception as e:
        st.error(f"❌ Error loading chronic issues data: {str(e)}")
        st.info("Please check your database connection and ensure the 'quality.clean_quality_data' table exists.")

def render_operator_trends(engine):
    """Show operator performance trends over time using monthly aggregation"""
    st.markdown("### 👥 Operator Performance Trends")
    st.info("Monthly defect and scrap analysis by operator (last 12 months)")

    # ---------------- SQL QUERY ----------------
    operator_query = """
    SELECT
        DATE_TRUNC('month', date)::date AS month,
        who_made_it AS operator_id,
        COUNT(*) AS defect_count,
        SUM(CASE WHEN disposition = 'SCRAP' THEN 1 ELSE 0 END) AS scrap_count
    FROM quality.clean_quality_data
    WHERE date >= CURRENT_DATE - INTERVAL '12 months'
      AND who_made_it IS NOT NULL
    GROUP BY DATE_TRUNC('month', date)::date, who_made_it
    ORDER BY month, who_made_it
    """

    # Load data from the database
    operator_data = pd.read_sql(operator_query, engine)

    if operator_data.empty:
        st.warning("No operator trend data available.")
        return

    # ---------------- DATA CLEANING ----------------
    operator_data['month'] = pd.to_datetime(operator_data['month'])
    operator_data['scrap_rate'] = np.where(
        operator_data['defect_count'] == 0,
        0,
        (operator_data['scrap_count'] / operator_data['defect_count']) * 100
    ).round(2)

    # ---------------- TOP OPERATORS ----------------
    top_operators = (
        operator_data.groupby('operator_id')['defect_count']
        .sum()
        .nlargest(10)
        .index
    )

    st.markdown("#### 🔝 Top Operators by Total Defects (Last 12 Months)")
    st.dataframe(
        operator_data.groupby('operator_id')['defect_count']
        .sum()
        .sort_values(ascending=False)
        .reset_index()
        .rename(columns={'defect_count': 'Total Defects'})
    )

    selected_operators = st.multiselect(
        "Select operators to analyze:",
        options=operator_data['operator_id'].unique(),
        default=list(top_operators[:5])
    )

    if not selected_operators:
        st.warning("Please select at least one operator.")
        return

    filtered = operator_data[operator_data['operator_id'].isin(selected_operators)]

    # ---------------- MONTHLY DEFECT TREND ----------------
    st.subheader("📈 Monthly Defects per Operator")

    defects_pivot = filtered.pivot_table(
        index='month',
        columns='operator_id',
        values='defect_count',
        aggfunc='sum'
    ).fillna(0)

    st.line_chart(defects_pivot, height=400)

    # ---------------- MONTHLY SCRAP RATE TREND ----------------
    st.subheader("📉 Monthly Scrap Rate (%) per Operator")

    scrap_pivot = filtered.pivot_table(
        index='month',
        columns='operator_id',
        values='scrap_rate',
        aggfunc='mean'
    ).fillna(0)

    st.line_chart(scrap_pivot, height=400)

    # ---------------- KPI SUMMARY ----------------
    st.subheader("📊 Operator Summary (Last 12 Months)")

    kpi = (
        filtered.groupby('operator_id')[['defect_count', 'scrap_count']]
        .sum()
        .reset_index()
    )
    kpi['scrap_rate (%)'] = (kpi['scrap_count'] / kpi['defect_count'] * 100).round(1)
    st.dataframe(kpi.sort_values('defect_count', ascending=False))

def render_performance_trends(engine):
    """Dynamic performance trends dashboard (Daily, Weekly, Monthly)"""
    st.markdown("## 📊 Performance Trends")
    st.info("Analyze defects and scrap rate trends by time period")

    # --- Time granularity selection ---
    view = st.selectbox(
        "Select time granularity:",
        ["Daily", "Weekly", "Monthly"],
        index=2  # Default: Monthly
    )

    # --- Build SQL query based on granularity ---
    if view == "Daily":
        query = """
        SELECT 
            DATE(date) AS period,
            COUNT(*) AS total_defects,
            COUNT(CASE WHEN UPPER(disposition) = 'SCRAP' THEN 1 END) AS scrap_count
        FROM quality.clean_quality_data
        WHERE date >= CURRENT_DATE - INTERVAL '30 days'
        GROUP BY DATE(date)
        ORDER BY period ASC
        """
    elif view == "Weekly":
        query = """
        SELECT 
            DATE_TRUNC('week', date)::date AS period,
            COUNT(*) AS total_defects,
            COUNT(CASE WHEN UPPER(disposition) = 'SCRAP' THEN 1 END) AS scrap_count
        FROM quality.clean_quality_data
        WHERE date >= CURRENT_DATE - INTERVAL '12 weeks'
        GROUP BY DATE_TRUNC('week', date)::date
        ORDER BY period ASC
        """
    else:  # Monthly
        query = """
        SELECT 
            DATE_TRUNC('month', date)::date AS period,
            COUNT(*) AS total_defects,
            COUNT(CASE WHEN UPPER(disposition) = 'SCRAP' THEN 1 END) AS scrap_count
        FROM quality.clean_quality_data
        WHERE date >= CURRENT_DATE - INTERVAL '12 months'
        GROUP BY DATE_TRUNC('month', date)::date
        ORDER BY period ASC
        """

    try:
        # --- Load data ---
        df = pd.read_sql(query, engine)
        if df.empty:
            st.warning(f"No {view.lower()} data available for analysis.")
            return

        # --- Prepare data ---
        df['period'] = pd.to_datetime(df['period'])
        df = df.sort_values('period', ascending=True)
        
        # Calculate scrap rate as percentage
        df['scrap_rate'] = (df['scrap_count'] / df['total_defects']) * 100
        df['scrap_rate'] = df['scrap_rate'].round(2)

        # --- Compute averages for comparison ---
        if len(df) > 1:
            latest = df.iloc[-1]
            avg_prev = df.iloc[:-1].mean(numeric_only=True)
        else:
            latest = df.iloc[0]
            avg_prev = latest

        col1, col2, col3, col4 = st.columns(4)
        with col1:
            delta = latest['total_defects'] - avg_prev['total_defects'] if len(df) > 1 else 0
            st.metric("Defects", f"{int(latest['total_defects'])}", f"{delta:+.0f} vs avg")
        with col2:
            delta = latest['scrap_count'] - avg_prev['scrap_count'] if len(df) > 1 else 0
            st.metric("Scrap", f"{int(latest['scrap_count'])}", f"{delta:+.0f} vs avg")
        with col3:
            delta = latest['scrap_rate'] - avg_prev['scrap_rate'] if len(df) > 1 else 0
            st.metric("Scrap Rate", f"{latest['scrap_rate']:.1f}%", f"{delta:+.1f}% vs avg")
        with col4:
            trend = "📈 Worse" if delta > 0 else "📉 Better" if len(df) > 1 else "➡️ Stable"
            st.metric("Trend Direction", trend)

        # --- ENHANCED INTERACTIVE CHARTS WITH ALT AIR ---
        import altair as alt
        
        # Format period for display
        if view == "Daily":
            df['period_display'] = df['period'].dt.strftime('%Y-%m-%d')
        elif view == "Weekly":
            df['period_display'] = df['period'].dt.strftime('Week of %Y-%m-%d')
        else:  # Monthly
            df['period_display'] = df['period'].dt.strftime('%b %Y')



        # Chart 2: Scrap Rate with enhanced tooltips
        st.subheader("📉 Scrap Rate Trend (%)")
        
        scrap_chart = alt.Chart(df).mark_line(point=True, strokeWidth=3, color='red').encode(
            x=alt.X('period:T', title=view, axis=alt.Axis(format='%b %d' if view == 'Daily' else '%b %Y')),
            y=alt.Y('scrap_rate:Q', title='Scrap Rate (%)', scale=alt.Scale(zero=True)),
            tooltip=[
                alt.Tooltip('period_display:N', title='Date'),
                alt.Tooltip('total_defects:Q', title='Defect Count', format=',.0f'),
                alt.Tooltip('scrap_count:Q', title='Scrap Count', format=',.0f'),
                alt.Tooltip('scrap_rate:Q', title='Scrap Rate %', format='.1f'),
                alt.Tooltip('scrap_count:Q', title='Scrap/Total', format=',.0f'),
                alt.Tooltip('total_defects:Q', title='Total Defects', format=',.0f')
            ]
        ).properties(
            height=400,
            title=f"{view} Scrap Rate Trend"
        ).interactive()
        
        st.altair_chart(scrap_chart, use_container_width=True)

        # --- Table view ---
        st.markdown(f"#### 📅 {view} Summary Table")
        display_df = df.rename(columns={
            'period': view,
            'total_defects': 'Total Defects',
            'scrap_count': 'Scrap Count',
            'scrap_rate': 'Scrap Rate (%)'
        }).copy()
        
        # Format the period column for display
        if view == "Daily":
            display_df[view] = display_df[view].dt.strftime('%Y-%m-%d')
        elif view == "Weekly":
            display_df[view] = display_df[view].dt.strftime('Week of %Y-%m-%d')
        else:  # Monthly
            display_df[view] = display_df[view].dt.strftime('%b %Y')
            
        st.dataframe(display_df, use_container_width=True)

    except Exception as e:
        st.error(f"Error loading {view.lower()} performance data: {e}")

def render_advanced_analysis(engine):
    """New advanced analyses from SQL queries"""
    st.markdown("### 🔍 Advanced Quality Analysis")
    
    # Analysis 1: Operator defect rate with machine details
    st.markdown("#### 👥 Operator-Machine Defect Analysis")
    operator_query = """
    SELECT 
        who_made_it as operator_id,
        code_description as defect,
        machine_no,
        COUNT(*) as defect_count,
        COUNT(CASE WHEN disposition = 'Scrap' THEN 1 END) as scrap_count
    FROM quality.clean_quality_data  
    WHERE date >= CURRENT_DATE - INTERVAL '30 days'
    GROUP BY 1, 2, 3
    ORDER BY 4 DESC
    LIMIT 50
    """
    
    operator_data = run_query(operator_query, engine)
    
    if not operator_data.empty:
        # Pivot table for heatmap
        pivot_data = operator_data.pivot_table(
            index='operator_id',
            columns='machine_no',
            values='defect_count',
            aggfunc='sum'
        ).fillna(0)
        
        if not pivot_data.empty and len(pivot_data) > 1:
            fig = px.imshow(
                pivot_data,
                title="Operator Defects by Machine (Heatmap)",
                aspect="auto",
                color_continuous_scale="reds"
            )
            st.plotly_chart(fig, use_container_width=True)
        
        # Top operator-machine combinations
        st.markdown("**Top Operator-Machine Defect Combinations**")
        display_data = operator_data.head(15)[['operator_id', 'machine_no', 'defect', 'defect_count', 'scrap_count']]
        display_data['scrap_rate'] = (display_data['scrap_count'] / display_data['defect_count'] * 100).round(1)
        st.dataframe(display_data, use_container_width=True)
    
    # Analysis 2: Top defective machines
    st.markdown("#### 🏭 Top Defective Machines Analysis")
    machine_query = """
    SELECT 
        machine_no,
        code_description as defect,
        COUNT(*) as defect_count
    FROM quality.clean_quality_data  
    WHERE date >= CURRENT_DATE - INTERVAL '30 days'
    GROUP BY 1, 2
    ORDER BY 3 DESC
    LIMIT 15
    """
    
    machine_data = run_query(machine_query, engine)
    
    if not machine_data.empty:
        fig = px.sunburst(
            machine_data,
            path=['machine_no', 'defect'],
            values='defect_count',
            title="Machine-Defect Relationship (Sunburst Chart)"
        )
        st.plotly_chart(fig, use_container_width=True)
    
    # Analysis 3: Monthly defect trends by category
    st.markdown("#### 📅 Monthly Defect Trends by Category")
    monthly_query = """
    SELECT 
        DATE_TRUNC('month', date) as month,
        category,
        COUNT(*) as defect_count
    FROM quality.clean_quality_data  
    WHERE date >= CURRENT_DATE - INTERVAL '6 months'
    GROUP BY 1, 2
    ORDER BY 1, 3 DESC
    """
    
    monthly_data = run_query(monthly_query, engine)
    
    if not monthly_data.empty:
        monthly_data['month'] = pd.to_datetime(monthly_data['month'])
        
        fig = px.line(
            monthly_data,
            x='month',
            y='defect_count',
            color='category',
            title="Monthly Defect Trends by Category",
            markers=True
        )
        st.plotly_chart(fig, use_container_width=True)
    
    # Analysis 4: Monthly operator performance
    st.markdown("#### 📊 Monthly Operator Performance")
    operator_monthly_query = """
    SELECT
        DATE_TRUNC('month', date) as month,
        who_made_it as operator_id,
        COUNT(*) as defect_count,
        COUNT(CASE WHEN disposition = 'Scrap' THEN 1 END) as scrap_count
    FROM quality.clean_quality_data
    WHERE date >= CURRENT_DATE - INTERVAL '6 months'
    GROUP BY 1, 2
    HAVING COUNT(*) >= 5
    ORDER BY 1, 3 DESC
    """
    
    operator_monthly_data = run_query(operator_monthly_query, engine)
    
    if not operator_monthly_data.empty:
        operator_monthly_data['month'] = pd.to_datetime(operator_monthly_data['month'])
        operator_monthly_data['scrap_rate'] = (operator_monthly_data['scrap_count'] / operator_monthly_data['defect_count'] * 100).round(1)
        
        # Top 10 operators by defect count
        top_operators = operator_monthly_data.groupby('operator_id')['defect_count'].sum().nlargest(10).index
        filtered_data = operator_monthly_data[operator_monthly_data['operator_id'].isin(top_operators)]
        
        fig = px.line(
            filtered_data,
            x='month',
            y='defect_count',
            color='operator_id',
            title="Top 10 Operators - Monthly Defect Trends",
            markers=True
        )
        st.plotly_chart(fig, use_container_width=True)

def create_modern_pareto_chart(series, title, xaxis_title, top_n=15):
    """Create a modern Pareto chart"""
    counts = series.value_counts().head(top_n)
    
    if counts.empty:
        return None, pd.DataFrame()
    
    # Calculate percentages and cumulative percentages
    total = counts.sum()
    percentages = (counts / total * 100).round(1)
    cumulative_percentages = percentages.cumsum()
    
    # Create figure with secondary y-axis
    fig = go.Figure()
    
    # Bar chart for counts
    fig.add_trace(go.Bar(
        x=counts.index,
        y=counts.values,
        name="Count",
        marker_color='#3366cc',
        text=counts.values,
        textposition='auto',
    ))
    
    # Line chart for cumulative percentage
    fig.add_trace(go.Scatter(
        x=counts.index,
        y=cumulative_percentages.values,
        name="Cumulative %",
        yaxis="y2",
        line=dict(color='#ff9900', width=3),
        marker=dict(size=8)
    ))
    
    fig.update_layout(
        title=title,
        xaxis_title=xaxis_title,
        yaxis_title="Count",
        yaxis2=dict(
            title="Cumulative Percentage",
            overlaying="y",
            side="right",
            range=[0, 100]
        ),
        xaxis_tickangle=-45,
        hovermode="x unified",
        height=500
    )
    
    # Prepare pareto data
    pareto_data = pd.DataFrame({
        'category': counts.index,
        'count': counts.values,
        'percentage': percentages.values,
        'cumulative_percentage': cumulative_percentages.values
    })
    
    return fig, pareto_data