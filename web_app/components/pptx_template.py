"""
Pretty PPTX template helper for the Pareto dashboard.

Provide create_pretty_pptx(plots, tables, title, logo_path=None, brand_color='#2C3E50', accent_color='#FF9900', title_font='Calibri', body_font='Calibri')
Returns: bytes (PPTX) or None if python-pptx not available.

Usage:
  from web_app.components.pptx_template import create_pretty_pptx
  pptx_bytes = create_pretty_pptx(plots, tables, "Quality Pareto", logo_path="assets/logo.png", brand_color="#123456")

Notes:
  - Requires python-pptx and pillow for image insertion.
  - If logo_path is a URL, the helper will try to fetch it via requests (if available). If requests isn't available,
    pass a local path or bytes.
"""
import io
from datetime import datetime

import pandas as pd

# Guard imports to fail gracefully
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    RGBColor = None
    MSO_SHAPE = None
    HAS_PPTX = False

# optional requests to fetch remote logo
try:
    import requests
    HAS_REQUESTS = True
except Exception:
    HAS_REQUESTS = False

# pillow used indirectly by python-pptx for some image operations; we don't need to import directly here
try:
    from PIL import Image as PILImage  # noqa: F401
    HAS_PIL = True
except Exception:
    HAS_PIL = False


def _hex_to_rgb_tuple(hex_color: str):
    hex_color = (hex_color or "#2C3E50").lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join([c * 2 for c in hex_color])
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def _get_logo_bytes(logo_path_or_url):
    if not logo_path_or_url:
        return None
    # If it's a URL and requests available fetch it
    if isinstance(logo_path_or_url, str) and (logo_path_or_url.startswith("http://") or logo_path_or_url.startswith("https://")):
        if not HAS_REQUESTS:
            return None
        try:
            r = requests.get(logo_path_or_url, timeout=10)
            if r.status_code == 200:
                return r.content
            return None
        except Exception:
            return None
    # Otherwise treat as file path or raw bytes
    if isinstance(logo_path_or_url, (bytes, bytearray)):
        return bytes(logo_path_or_url)
    try:
        with open(logo_path_or_url, "rb") as f:
            return f.read()
    except Exception:
        return None


def create_pretty_pptx(
    plots: dict,
    tables: dict,
    title: str = "Quality Pareto Analysis",
    logo_path: str | None = None,
    brand_color: str = "#2C3E50",
    accent_color: str = "#FF9900",
    title_font: str = "Calibri",
    body_font: str = "Calibri",
) -> bytes | None:
    """
    Build a styled PPTX:
      - Title slide with optional logo
      - Executive summary slide (bullets)
      - One slide per plot (chart image embedded)
      - One slide per table (top rows as PPTX table)
    Returns bytes or None if prerequisites missing.
    """
    if not HAS_PPTX:
        return None

    # small defensive: convert colors to RGB
    brand_rgb = _hex_to_rgb_tuple(brand_color)
    accent_rgb = _hex_to_rgb_tuple(accent_color)

    prs = Presentation()
    # Title slide (large, centered)
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
        try:
            title_tf = slide.shapes.title.text_frame
            title_tf.paragraphs[0].font.name = title_font
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor(*brand_rgb)
        except Exception:
            pass

    # Add logo top-right if available
    logo_bytes = _get_logo_bytes(logo_path)
    if logo_bytes:
        try:
            logo_stream = io.BytesIO(logo_bytes)
            # place near top-right; width 1.6" (adjust as needed)
            slide.shapes.add_picture(logo_stream, prs.slide_width - Inches(2.0), Inches(0.4), width=Inches(1.6))
        except Exception:
            # ignore logo failure
            pass

    # Subtitle / timestamp
    try:
        # try placeholder 1
        subtitle = slide.placeholders[1]
        subtitle.text = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        for p in subtitle.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = body_font
            p.font.color.rgb = RGBColor(*accent_rgb)
    except Exception:
        # no placeholder available, continue
        pass

    # Executive summary slide
    summary_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[5]
    sslide = prs.slides.add_slide(summary_layout)
    try:
        if sslide.shapes.title:
            sslide.shapes.title.text = "Executive Summary"
            sslide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*brand_rgb)
    except Exception:
        pass
    tx = sslide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    # Build bullets intelligently from provided tables/plots
    bullets = []
    if "Top Defects" in tables and isinstance(tables["Top Defects"], pd.DataFrame) and not tables["Top Defects"].empty:
        t = tables["Top Defects"]
        bullets.append(f"Top defects listed: {len(t)}")
        if "count" in t.columns:
            bullets.append(f"Sum of top defect counts: {int(t['count'].sum()):,}")
        if "cumulative_percentage" in t.columns:
            bullets.append(f"Top cumulative % (first row): {t['cumulative_percentage'].iat[0]:.1f}%")
    else:
        bullets.append("No top-defect table included.")

    bullets.append(f"Charts included: {', '.join(list(plots.keys()))}" if plots else "No charts included.")
    bullets.append(f"Report generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.level = 0
        p.font.name = body_font
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(*brand_rgb)

    # Plot slides
    for p_title, fig in plots.items():
        s_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
        s = prs.slides.add_slide(s_layout)
        try:
            if s.shapes.title:
                s.shapes.title.text = p_title
                s.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*brand_rgb)
        except Exception:
            pass
        # Try to render chart to PNG bytes - prefer plotly.kaleido via fig.to_image
        try:
            png = fig.to_image(format="png", width=1400, height=700)
        except Exception:
            try:
                buf = io.BytesIO()
                fig.write_image(buf, format="png", width=1400, height=700)
                png = buf.getvalue()
            except Exception:
                png = None
        if png:
            try:
                img_buf = io.BytesIO(png)
                s.shapes.add_picture(img_buf, Inches(0.5), Inches(1.2), width=Inches(9.0))
            except Exception:
                # fallback to text if insertion fails
                tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.0))
                tb.text = "Chart image unavailable (ensure kaleido installed)."
        else:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.0))
            tb.text = "Chart image unavailable (ensure kaleido installed)."

    # Table slides: create a slide with a formatted PPTX table for each table (top up to 10 rows)
    for t_title, table in tables.items():
        if not isinstance(table, pd.DataFrame) or table.empty:
            continue
        tbl = table.copy()
        rows = min(10, len(tbl))
        cols = len(tbl.columns)
        # create slide
        s_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
        s = prs.slides.add_slide(s_layout)
        try:
            if s.shapes.title:
                s.shapes.title.text = t_title
                s.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*brand_rgb)
        except Exception:
            pass
        # Add table with header row + rows rows
        try:
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(9.0)
            height = Inches(5.0)
            table_shape = s.shapes.add_table(rows + 1, cols, left, top, width, height).table
            # Header
            for c, col_name in enumerate(tbl.columns):
                cell = table_shape.cell(0, c)
                cell.text = str(col_name)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = body_font
                    paragraph.font.color.rgb = RGBColor(*brand_rgb)
            # Data rows
            for r in range(rows):
                for c, col_name in enumerate(tbl.columns):
                    cell = table_shape.cell(r + 1, c)
                    val = tbl.iloc[r, c]
                    cell.text = str(val)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.name = body_font
        except Exception:
            # fallback to bullet text
            tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.0))
            tf = tb.text_frame
            for r in tbl.head(rows).to_dict("records"):
                line = " • ".join([f"{k}: {v}" for k, v in r.items()])
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)

    # Finalize
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()